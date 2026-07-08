"""结构化办公文档 → Markdown（Excel / PPTX）。

表格统一渲染为「字段名：内容 字段名：内容」纯文本（每行一条 record），
不再输出 ``| ... |`` MD 表格 —— 便于知识库切片与向量检索。详见
:mod:`file_to_markdown.table_renderer`。
"""

from __future__ import annotations

import io
from typing import Any

from .table_renderer import table_to_field_value_text


def excel_bytes_to_markdown(content: bytes) -> tuple[str, dict[str, Any]]:
    """
    Excel → Markdown。

    每个 Sheet 一段 ``## Sheet: <name>`` 标题 + 「字段名：内容 …」每行一条 record 的纯文本。
    多 Sheet 之间用 ``---`` 分隔。Sheet 内无有效行时仅保留标题。
    """
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheet_names = list(wb.sheetnames)
    sheets_md: list[str] = []
    record_count = 0

    try:
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            sheet_parts: list[str] = [f"## Sheet: {sheet_name}\n"]

            rows_raw: list[list[Any]] = []
            for row in ws.iter_rows(values_only=True):
                if not any(
                    (c is not None and str(c).strip() != "") for c in row
                ):
                    continue
                rows_raw.append(list(row))

            if rows_raw:
                headers = rows_raw[0]
                data_rows = rows_raw[1:]
                text = table_to_field_value_text(headers, data_rows)
                if text:
                    sheet_parts.append(text)
                    record_count += text.count("\n") + 1

            sheets_md.append("\n\n".join(sheet_parts))
    finally:
        wb.close()

    md = "\n\n---\n\n".join(sheets_md) if sheets_md else ""
    return md, {
        "sheet_count": len(sheet_names),
        "record_count": record_count,
        "converter": "excel_field_value",
        "table_format": "field_value_text",
    }


def pptx_bytes_to_markdown(content: bytes) -> tuple[str, dict[str, Any]]:
    """
    PPT 文本版 → Markdown。

    - 每页 ``## Slide N`` 标题 + 文本框正文；
    - 页内表格也走 ``table_to_field_value_text``（每行一条 record 的「字段名：内容 ...」纯文本），
      不再输出 ``| ... |`` MD 表格。
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides_md: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = [f"## Slide {slide_num}\n"]
        texts: list[str] = []

        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = (para.text or "").strip()
                    if t:
                        texts.append(t)

            if getattr(shape, "has_table", False) and shape.has_table:
                table = shape.table
                rows_raw: list[list[str]] = []
                for row in table.rows:
                    cells = [(cell.text or "") for cell in row.cells]
                    if any(c.strip() for c in cells):
                        rows_raw.append(cells)
                if len(rows_raw) >= 1:
                    headers = rows_raw[0]
                    data_rows = rows_raw[1:]
                    table_text = table_to_field_value_text(headers, data_rows)
                    if table_text:
                        texts.append(table_text)

        if texts:
            slide_parts.append("\n\n".join(texts))
        slides_md.append("\n\n".join(slide_parts))

    md = "\n\n---\n\n".join(slides_md)
    return md, {
        "slide_count": len(prs.slides),
        "converter": "pptx_structured",
        "table_format": "field_value_text",
    }
